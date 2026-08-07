import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

print("Opening presentation...")
p = Presentation(r'e:\CAP\RestockIQ_Review-I_MASTER.pptx')
slide = p.slides[5] # Slide 6 is index 5

print("Removing old shapes...")
# Remove all shapes except the title (assuming title is the first shape or has text 'Objectives of the Project')
shapes_to_remove = []
for sh in slide.shapes:
    if not sh.has_text_frame or 'Objectives' not in sh.text_frame.text:
        shapes_to_remove.append(sh)

for sh in shapes_to_remove:
    sh._element.getparent().remove(sh._element)

print("Adding styling and content...")
# Add a thin separator line under the title (matches deck style)
line = slide.shapes.add_connector(1, Emu(382554), Emu(1500000), Emu(382554)+Emu(11155680-382554*2), Emu(1500000))
line.line.color.rgb = RGBColor(0xE0, 0xDA, 0xD8)
line.line.width = Pt(1)

MAROON = RGBColor(0xBC, 0x31, 0x2C)
NAVY = RGBColor(0x2F, 0x3C, 0x7E)
CARDBG = RGBColor(0xF7, 0xF5, 0xF4)
BORDER = RGBColor(0xE0, 0xDA, 0xD8)
BODY = RGBColor(0x33, 0x33, 0x33)

objectives = [
    {
        "num": "01",
        "title": "Build the Forecasting Model",
        "subs": [
            "Design an adaptive Prophet + XGBoost hybrid \u2014 Prophet for trend/seasonality, XGBoost for residual correction \u2014 with regime-specific tuning for sparse and high-volume items."
        ],
    },
    {
        "num": "02",
        "title": "Generate Item-Level Forecasts",
        "subs": [
            "Produce daily, item-level demand predictions from the trained hybrid model, validated for consistency across sparse and high-volume item categories."
        ],
    },
    {
        "num": "03",
        "title": "Recommend Restock Quantities",
        "subs": [
            "Convert raw forecasts into actionable restock quantities using a configurable safety-stock buffer that accounts for current inventory levels."
        ],
    },
    {
        "num": "04",
        "title": "Explain via LLM Narrator",
        "subs": [
            "Integrate an LLM layer that translates the restock output into a plain-language recommendation \u2014 understandable without technical or data literacy."
        ],
    },
]

# 1x4 column grid geometry (EMU). Slide width ~12192000 (13.33in), height 6858000 (7.5in) typical 16:9
slide_w = p.slide_width
slide_h = p.slide_height

margin_x = Emu(400000)
gap_x = Emu(250000)
top = Emu(1800000)
# calculate width for 4 columns
card_w = Emu(int((slide_w - 2*margin_x - 3*gap_x) / 4))
card_h = Emu(4200000)

positions = [
    (margin_x, top),
    (Emu(margin_x + card_w + gap_x), top),
    (Emu(margin_x + 2*card_w + 2*gap_x), top),
    (Emu(margin_x + 3*card_w + 3*gap_x), top),
]

for obj, (x, y) in zip(objectives, positions):
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARDBG
    card.line.color.rgb = BORDER
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # Number
    num_box = slide.shapes.add_textbox(Emu(x+150000), Emu(y+100000), Emu(card_w-300000), Emu(500000))
    tf = num_box.text_frame
    tf.margin_left = 0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = obj["num"]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.name = "Cambria"
    run.font.color.rgb = MAROON

    # Title
    title_box = slide.shapes.add_textbox(Emu(x+150000), Emu(y+700000), Emu(card_w-300000), Emu(800000))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = obj["title"]
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = NAVY

    # Sub-objectives (no bullet, just text per the image)
    body_box = slide.shapes.add_textbox(Emu(x+150000), Emu(y+1600000), Emu(card_w-300000), Emu(card_h-1800000))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, sub in enumerate(obj["subs"]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = sub
        run.font.size = Pt(14)
        run.font.name = "Calibri"
        run.font.color.rgb = BODY

output_file = r'e:\CAP\RestockIQ_Review-I_MASTER_updated_v2.pptx'
p.save(output_file)
print(f"Saved to {output_file}")
