import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

print("Opening presentation...")
# Let's open the original file or the v3 one if it exists. We'll use the original and save a new one.
p = Presentation(r'e:\CAP\RestockIQ_Review-I_MASTER.pptx')
slide = p.slides[9] # Slide 10 is index 9

print("Removing old shapes...")
shapes_to_remove = []
for sh in slide.shapes:
    if not sh.has_text_frame or 'Tools' not in sh.text_frame.text:
        shapes_to_remove.append(sh)

for sh in shapes_to_remove:
    sh._element.getparent().remove(sh._element)

print("Adding styling and content...")
line = slide.shapes.add_connector(1, Emu(382554), Emu(1400000), Emu(382554)+Emu(11155680-382554*2), Emu(1400000))
line.line.color.rgb = RGBColor(0xE0, 0xDA, 0xD8)
line.line.width = Pt(1)

RED_ACCENT = RGBColor(0xBC, 0x31, 0x2C)
NAVY = RGBColor(0x2F, 0x3C, 0x7E)
BODY = RGBColor(0x33, 0x33, 0x33)

cards_data = [
    {
        "title": "Mobile Frontend",
        "subtitle": "Flutter (Dart)",
        "text": "Cross-platform app; zero-install-friction interface for non-technical owners"
    },
    {
        "title": "Forecasting Engine",
        "subtitle": "Prophet + XGBoost",
        "text": "Prophet handles trend/seasonality; XGBoost models the residuals."
    },
    {
        "title": "External Regressors",
        "subtitle": "WeatherAPI, `holidays`",
        "text": "Feeds weather and public-holiday signals into the Prophet model."
    },
    {
        "title": "Restock Logic",
        "subtitle": "Python (Pandas, NumPy)",
        "text": "Deterministic safety-stock buffer calculation \u2014 not a second ML model."
    },
    {
        "title": "Backend API",
        "subtitle": "FastAPI / Flask",
        "text": "Connects Flutter app to the Python forecasting pipeline."
    },
    {
        "title": "Database",
        "subtitle": "Firebase / Supabase",
        "text": "User auth, shop profiles, and historical sales logs."
    },
    {
        "title": "Development Tools",
        "subtitle": "Jupyter, VS Code, Git",
        "text": "Jupyter Notebook for EDA & model training; VS Code for app/API development."
    },
    {
        "title": "Hardware Specs",
        "subtitle": "Standard CPU / 8GB RAM",
        "text": "No GPU required. Computationally lightweight, ensuring low cloud-inference costs."
    }
]

# 2x4 grid (4 columns, 2 rows)
slide_w = p.slide_width
slide_h = p.slide_height

margin_x = Emu(400000)
gap_x = Emu(200000)
gap_y = Emu(300000)
top = Emu(1700000)

card_w = Emu(int((slide_w - 2*margin_x - 3*gap_x)/4))
card_h = Emu(int((slide_h - top - gap_y - Emu(400000))/2))

positions = []
for row in range(2):
    for col in range(4):
        x = margin_x + col * (card_w + gap_x)
        y = top + row * (card_h + gap_y)
        positions.append((Emu(x), Emu(y)))

for data, (x, y) in zip(cards_data, positions):
    # Left border red accent line
    accent = slide.shapes.add_connector(1, x, y, x, y+card_h)
    accent.line.color.rgb = RED_ACCENT
    accent.line.width = Pt(4)

    # Box frame (optional, matches the subtle grey box in original)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
    card.fill.background()
    card.line.color.rgb = RGBColor(0xE0, 0xDA, 0xD8)
    card.line.width = Pt(1)

    # Text Box
    tb = slide.shapes.add_textbox(x + Emu(100000), y + Emu(100000), card_w - Emu(200000), card_h - Emu(200000))
    tf = tb.text_frame
    tf.word_wrap = True
    
    # Title
    p_title = tf.paragraphs[0]
    r_title = p_title.add_run()
    r_title.text = data["title"]
    r_title.font.name = "Calibri"
    r_title.font.size = Pt(16)
    r_title.font.bold = True
    r_title.font.color.rgb = NAVY
    
    # Subtitle
    p_sub = tf.add_paragraph()
    p_sub.space_before = Pt(10)
    r_sub = p_sub.add_run()
    r_sub.text = data["subtitle"]
    r_sub.font.name = "Calibri"
    r_sub.font.size = Pt(12)
    r_sub.font.bold = True
    r_sub.font.color.rgb = RED_ACCENT
    
    # Body
    p_body = tf.add_paragraph()
    p_body.space_before = Pt(10)
    r_body = p_body.add_run()
    r_body.text = data["text"]
    r_body.font.name = "Calibri"
    r_body.font.size = Pt(12)
    r_body.font.color.rgb = BODY

output_file = r'e:\CAP\RestockIQ_Review-I_MASTER_updated_tools.pptx'
p.save(output_file)
print(f"Saved to {output_file}")
