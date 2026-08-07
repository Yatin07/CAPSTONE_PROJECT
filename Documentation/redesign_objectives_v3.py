import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

print("Opening presentation...")
p = Presentation(r'e:\CAP\RestockIQ_Review-I_MASTER.pptx')
slide = p.slides[5] # Slide 6 is index 5

print("Removing old shapes...")
shapes_to_remove = []
for sh in slide.shapes:
    if not sh.has_text_frame or 'Objectives' not in sh.text_frame.text:
        shapes_to_remove.append(sh)

for sh in shapes_to_remove:
    sh._element.getparent().remove(sh._element)

print("Adding styling and content...")
line = slide.shapes.add_connector(1, Emu(382554), Emu(1500000), Emu(382554)+Emu(11155680-382554*2), Emu(1500000))
line.line.color.rgb = RGBColor(0xE0, 0xDA, 0xD8)
line.line.width = Pt(1)

MAROON = RGBColor(0xBC, 0x31, 0x2C)
NAVY = RGBColor(0x2F, 0x3C, 0x7E)
CARDBG = RGBColor(0xF7, 0xF5, 0xF4)
BORDER = RGBColor(0xE0, 0xDA, 0xD8)
BODY = RGBColor(0x33, 0x33, 0x33)

objectives = [
    {
        "num": "01",
        "title": "Build the Forecasting Model",
        "text": "Design an adaptive Prophet + XGBoost hybrid \u2014 Prophet for trend/seasonality, XGBoost for residual correction \u2014 with regime-specific tuning for sparse and high-volume items."
    },
    {
        "num": "02",
        "title": "Generate Item-Level Forecasts",
        "text": "Produce daily, item-level demand predictions from the trained hybrid model, validated for consistency across sparse and high-volume item categories."
    },
    {
        "num": "03",
        "title": "Recommend Restock Quantities",
        "text": "Convert raw forecasts into actionable restock quantities using a configurable safety-stock buffer that accounts for current inventory levels."
    },
    {
        "num": "04",
        "title": "Explain via LLM Narrator",
        "text": "Integrate an LLM layer that translates the restock output into a plain-language recommendation \u2014 understandable without technical or data literacy."
    },
]

# 2x2 grid geometry
slide_w = p.slide_width
slide_h = p.slide_height

margin_x = Emu(400000)
gap_x = Emu(300000)
gap_y = Emu(300000)
top = Emu(1800000)
card_w = Emu(int((slide_w - 2*margin_x - gap_x)//2))
card_h = Emu(2200000)

positions = [
    (margin_x, top),
    (Emu(margin_x + card_w + gap_x), top),
    (margin_x, Emu(top + card_h + gap_y)),
    (Emu(margin_x + card_w + gap_x), Emu(top + card_h + gap_y)),
]

for obj, (x, y) in zip(objectives, positions):
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARDBG
    card.line.color.rgb = BORDER
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # Number
    num_box = slide.shapes.add_textbox(Emu(x+150000), Emu(y+100000), Emu(900000), Emu(500000))
    tf = num_box.text_frame
    tf.margin_left = 0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = obj["num"]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.name = "Cambria"
    run.font.color.rgb = MAROON

    # Title
    title_box = slide.shapes.add_textbox(Emu(x+150000), Emu(y+560000), Emu(card_w-300000), Emu(430000))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = obj["title"]
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = NAVY

    # Body text
    body_box = slide.shapes.add_textbox(Emu(x+150000), Emu(y+1020000), Emu(card_w-300000), Emu(card_h-1120000))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = obj["text"]
    run.font.size = Pt(14)
    run.font.name = "Calibri"
    run.font.color.rgb = BODY

output_file = r'e:\CAP\RestockIQ_Review-I_MASTER_updated_v3.pptx'
p.save(output_file)
print(f"Saved to {output_file}")
